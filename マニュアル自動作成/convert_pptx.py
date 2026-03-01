import os
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
import html
import shutil

# Files
PPTX_FILE = "C:/Users/24030460/Desktop/antigravity/マニュアル/20230509_マニュアルフォーマットver1.pptx"
OUTPUT_DIR = "C:/Users/24030460/Desktop/antigravity/マニュアル/manual_output"
ASSETS_DIR = os.path.join(OUTPUT_DIR, "assets")
HTML_FILE = os.path.join(OUTPUT_DIR, "manual_view.html")

# Ensure output directories exist
if os.path.exists(OUTPUT_DIR):
    shutil.rmtree(OUTPUT_DIR)
os.makedirs(ASSETS_DIR)

print(f"Loading {PPTX_FILE}...")
prs = pptx.Presentation(PPTX_FILE)

# Slide dimensions (EMU)
SLIDE_WIDTH_EMU = prs.slide_width
SLIDE_HEIGHT_EMU = prs.slide_height

# Scale factor to pixels (approx 96 DPI)
EMU_PER_INCH = 914400
PIXELS_PER_INCH = 96
SCALE = PIXELS_PER_INCH / EMU_PER_INCH

SLIDE_WIDTH_PX = int(SLIDE_WIDTH_EMU * SCALE)
SLIDE_HEIGHT_PX = int(SLIDE_HEIGHT_EMU * SCALE)

html_content = f"""<!DOCTYPE html>
<html lang="ja">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Manual Viewer</title>
    <style>
        body {{
            background-color: #f0f0f0;
            font-family: sans-serif;
            margin: 0;
            padding: 20px;
            display: flex;
            flex-direction: column;
            align-items: center;
        }}
        .slide {{
            position: relative;
            width: {SLIDE_WIDTH_PX}px;
            height: {SLIDE_HEIGHT_PX}px;
            background-color: white;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            margin-bottom: 20px;
            overflow: hidden;
        }}
        .element {{
            position: absolute;
            overflow: hidden;
            word-wrap: break-word;
        }}
        .slide-number {{
            position: absolute;
            bottom: 10px;
            right: 10px;
            font-size: 12px;
            color: #888;
        }}
    </style>
</head>
<body>
"""

def extract_shape(shape, slide_idx, shape_idx):
    style = ""
    content = ""
    
    # Position
    left = int(shape.left * SCALE)
    top = int(shape.top * SCALE)
    width = int(shape.width * SCALE)
    height = int(shape.height * SCALE)
    
    style += f"left: {left}px; top: {top}px; width: {width}px; height: {height}px;"

    # Text
    if shape.has_text_frame:
        text_html = ""
        for paragraph in shape.text_frame.paragraphs:
            p_style = "margin: 0;"
            # Alignment
            if paragraph.alignment:
                align_map = {1: 'left', 2: 'center', 3: 'right', 4: 'justify'}
                try:
                    val = int(paragraph.alignment)
                    if val in align_map:
                        p_style += f"text-align: {align_map[val]};"
                except:
                    pass
 
            
            p_text = ""
            for run in paragraph.runs:
                run_text = html.escape(run.text)
                run_style = ""
                if run.font.bold:
                    run_style += "font-weight: bold;"
                if run.font.italic:
                    run_style += "font-style: italic;"
                if run.font.size:
                    font_px = int(run.font.size.pt * 1.33) # pt to px approx
                    run_style += f"font-size: {font_px}px;"
                if run.font.color and hasattr(run.font.color, 'rgb'):
                     if run.font.color.rgb:
                        run_style += f"color: #{run.font.color.rgb};"
                
                p_text += f'<span style="{run_style}">{run_text}</span>'
            
            if not p_text.strip():
                p_text = "&nbsp;"
            text_html += f'<p style="{p_style}">{p_text}</p>'
        
        content = text_html
        # Check for fill
        if shape.fill.type == 1: # Solid check
             if shape.fill.fore_color.type == 1: # RGB
                 style += f"background-color: #{shape.fill.fore_color.rgb};"

    # Image
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        ext = image.ext
        img_filename = f"slide{slide_idx}_shape{shape_idx}.{ext}"
        img_path = os.path.join(ASSETS_DIR, img_filename)
        
        with open(img_path, "wb") as f:
            f.write(image.blob)
            
        content = f'<img src="assets/{img_filename}" style="width: 100%; height: 100%; object-fit: contain;">'

    # Group (Recursive - simple flattened approach for now or skip)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # Handling groups is complex in flat HTML, skipping for MVP or treating as container
        pass
        
    return f'<div class="element" style="{style}">{content}</div>'

for i, slide in enumerate(prs.slides):
    print(f"Processing slide {i+1}...")
    
    # Try to determine background
    bg_style = ""
    # Very basic background check - PPTX backgrounds are complex
    # Defaulting to white in CSS, but checking for slide-specific solid fills
    
    html_content += f'<div class="slide" id="slide{i+1}" style="{bg_style}">'
    
    for j, shape in enumerate(slide.shapes):
        try:
            html_content += extract_shape(shape, i+1, j)
        except Exception as e:
            import traceback
            print(f"Error processing shape {j} on slide {i+1}: {e}")
            traceback.print_exc()
            
    html_content += f'<div class="slide-number">{i+1}</div>'
    html_content += '</div>\n'

html_content += """
</body>
</html>
"""

with open(HTML_FILE, "w", encoding="utf-8") as f:
    f.write(html_content)

print(f"Done! HTML saved to {HTML_FILE}")
