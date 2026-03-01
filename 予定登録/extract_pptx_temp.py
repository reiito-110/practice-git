import pptx
from pptx import Presentation
import json
import os

def extract_pptx_data(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        return

    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_content = {
            "slide_number": i + 1,
            "title": "",
            "elements": []
        }
        
        if slide.shapes.title:
            slide_content["title"] = slide.shapes.title.text

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Skip title if we already got it, but sometimes titles are just text boxes
            if shape == slide.shapes.title:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    slide_content["elements"].append({
                        "text": text,
                        "level": paragraph.level
                    })
        
        slides_data.append(slide_content)

    return slides_data

if __name__ == "__main__":
    pptx_path = r"c:\Users\24030460\Desktop\antigravity\予定登録\20230509_マニュアルフォーマットver1.pptx"
    data = extract_pptx_data(pptx_path)
    print(json.dumps(data, ensure_ascii=False, indent=2))
